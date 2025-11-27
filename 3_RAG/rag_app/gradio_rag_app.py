"""
Gradio Web App for Document Indexing and RAG Queries with Qdrant

This app allows users to:
1. Upload multiple documents (PDF, Word, PPT, TXT)
2. Index them in Qdrant vector database
3. Ask questions about the indexed documents using RAG
"""

import os
import tempfile
from typing import List, Optional, Tuple
from pathlib import Path

import gradio as gr
from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_openai import OpenAIEmbeddings, ChatOpenAI
from langchain_community.vectorstores import Qdrant
from langchain_community.document_loaders import (
    PyPDFLoader,
    TextLoader,
)
from docx import Document as DocxDocument
from pptx import Presentation
from langchain_core.prompts import ChatPromptTemplate
from qdrant_client import QdrantClient
from qdrant_client.models import Distance, VectorParams

# Try to load environment variables
try:
    from dotenv import load_dotenv
    load_dotenv()
except ImportError:
    pass

# Global variables
embeddings: Optional[OpenAIEmbeddings] = None
vectorstore: Optional[Qdrant] = None
llm: Optional[ChatOpenAI] = None
retriever = None
collection_name = "gradio_rag_documents"

# Initialize components
def initialize_components():
    """Initialize embeddings, LLM, and Qdrant connection"""
    global embeddings, llm, vectorstore, retriever
    
    # Check for API key
    if not os.getenv("OPENAI_API_KEY"):
        raise ValueError("OPENAI_API_KEY not found in environment variables!")
    
    # Initialize embeddings and LLM
    embeddings = OpenAIEmbeddings(model="text-embedding-3-small")
    llm = ChatOpenAI(model="gpt-4o-mini", temperature=0)
    
    # Initialize Qdrant client (using in-memory for demo, can be changed to persistent)
    client = QdrantClient(location=":memory:")
    
    # Create collection if it doesn't exist
    try:
        client.get_collection(collection_name)
    except Exception:
        client.create_collection(
            collection_name=collection_name,
            vectors_config=VectorParams(
                size=1536,  # text-embedding-3-small dimension
                distance=Distance.COSINE
            )
        )
    
    # Create vectorstore
    vectorstore = Qdrant(
        client=client,
        collection_name=collection_name,
        embeddings=embeddings
    )
    
    # Create retriever with higher k to search across all documents
    # Using MMR (Maximal Marginal Relevance) for better diversity across documents
    try:
        retriever = vectorstore.as_retriever(
            search_type="mmr",  # Use MMR for diverse results
            search_kwargs={"k": 10, "fetch_k": 20}  # Retrieve more chunks to ensure coverage
        )
    except Exception:
        # Fallback to similarity search if MMR not supported
        retriever = vectorstore.as_retriever(
            search_kwargs={"k": 10}  # Still retrieve more chunks
        )
    
    return "✅ Components initialized successfully!"


def load_document(file_path: str) -> List[Document]:
    """Load a document based on its file extension"""
    file_ext = Path(file_path).suffix.lower()
    
    try:
        if file_ext == ".pdf":
            loader = PyPDFLoader(file_path)
            documents = loader.load()
        elif file_ext in [".doc", ".docx"]:
            # Load Word document using python-docx
            doc = DocxDocument(file_path)
            text_content = []
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_content.append(paragraph.text)
            # Also extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            text_content.append(cell.text)
            
            full_text = "\n".join(text_content)
            if not full_text.strip():
                raise ValueError("Document appears to be empty")
            
            documents = [Document(page_content=full_text, metadata={"source": Path(file_path).name})]
        elif file_ext in [".ppt", ".pptx"]:
            # Load PowerPoint document using python-pptx
            prs = Presentation(file_path)
            text_content = []
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)
                if slide_text:
                    text_content.append(f"Slide {slide_num}:\n" + "\n".join(slide_text))
            
            full_text = "\n\n".join(text_content)
            if not full_text.strip():
                raise ValueError("Presentation appears to be empty")
            
            documents = [Document(page_content=full_text, metadata={"source": Path(file_path).name})]
        elif file_ext == ".txt":
            loader = TextLoader(file_path, encoding="utf-8")
            documents = loader.load()
        else:
            raise ValueError(f"Unsupported file type: {file_ext}")
        
        return documents
    except Exception as e:
        raise Exception(f"Error loading document {file_path}: {str(e)}")


def index_documents(files: List[gr.File]) -> str:
    """Index uploaded documents to Qdrant"""
    global vectorstore, embeddings
    
    if not files:
        return "⚠️ Please upload at least one document!"
    
    if vectorstore is None:
        return "⚠️ Components not initialized! Please restart the app."
    
    all_documents = []
    processed_files = []
    errors = []
    
    # Process each uploaded file
    for file in files:
        if file is None:
            continue
            
        file_path = file.name if hasattr(file, 'name') else file
        
        try:
            # Load document
            docs = load_document(file_path)
            
            # Add metadata
            filename = Path(file_path).name
            for doc in docs:
                doc.metadata["source"] = filename
                doc.metadata["file_path"] = file_path
            
            all_documents.extend(docs)
            processed_files.append(filename)
            
        except Exception as e:
            errors.append(f"{Path(file_path).name}: {str(e)}")
    
    if not all_documents:
        error_msg = "⚠️ No documents could be loaded!\n"
        if errors:
            error_msg += "\nErrors:\n" + "\n".join(f"- {e}" for e in errors)
        return error_msg
    
    # Split documents into chunks with better overlap for context preservation
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=1000,
        chunk_overlap=200,
        length_function=len,
        separators=["\n\n", "\n", ". ", " ", ""]  # Better separators for natural breaks
    )
    chunks = text_splitter.split_documents(all_documents)
    
    # Add chunk index to metadata for better tracking
    for idx, chunk in enumerate(chunks):
        if "chunk_index" not in chunk.metadata:
            chunk.metadata["chunk_index"] = idx
    
    # Add to vectorstore
    try:
        vectorstore.add_documents(chunks)
        
        # Update retriever with MMR for diverse retrieval
        global retriever
        try:
            retriever = vectorstore.as_retriever(
                search_type="mmr",
                search_kwargs={"k": 10, "fetch_k": 20}
            )
        except Exception:
            # Fallback to similarity search if MMR not supported
            retriever = vectorstore.as_retriever(
                search_kwargs={"k": 10}
            )
        
        success_msg = f"✅ Successfully indexed {len(chunks)} chunks from {len(processed_files)} file(s)!\n\n"
        success_msg += f"Files processed:\n" + "\n".join(f"- {f}" for f in processed_files)
        
        if errors:
            success_msg += f"\n\n⚠️ Errors:\n" + "\n".join(f"- {e}" for e in errors)
        
        return success_msg
        
    except Exception as e:
        return f"❌ Error indexing documents: {str(e)}"


def query_rag(question: str, history: List[Tuple[str, str]]) -> Tuple[str, List[Tuple[str, str]]]:
    """Perform RAG query and return answer with detailed source attribution"""
    global retriever, llm
    
    if not question.strip():
        return "", history
    
    if retriever is None or llm is None:
        return "⚠️ Please initialize components and index documents first!", history
    
    try:
        # Retrieve relevant documents (MMR ensures diversity across all documents)
        docs = retriever.invoke(question)
        
        if not docs:
            return "⚠️ No relevant documents found. Please index some documents first.", history
        
        # Group documents by source for better organization
        sources_dict = {}
        for doc in docs:
            source = doc.metadata.get("source", "Unknown")
            if source not in sources_dict:
                sources_dict[source] = []
            sources_dict[source].append(doc)
        
        # Create context with clear source attribution
        context_parts = []
        for source, source_docs in sources_dict.items():
            for idx, doc in enumerate(source_docs, 1):
                chunk_preview = doc.page_content[:200] + "..." if len(doc.page_content) > 200 else doc.page_content
                context_parts.append(
                    f"[Document: {source} | Chunk {idx}]\n{doc.page_content}"
                )
        
        context = "\n\n---\n\n".join(context_parts)
        
        # Enhanced prompt with better instructions for source citation
        prompt = ChatPromptTemplate.from_messages([
            ("system", """You are a helpful assistant that answers questions based on the provided context from multiple documents.

IMPORTANT INSTRUCTIONS:
1. Synthesize information from ALL relevant documents provided in the context
2. If information appears in multiple documents, mention all relevant sources
3. Always cite the specific document name when referencing information
4. If the answer cannot be found in the context, explicitly state this
5. Be comprehensive and draw connections between information from different documents when relevant
6. Format your answer clearly with proper structure"""),
            ("human", """Context from indexed documents:

{context}

Question: {question}

Provide a comprehensive answer based on the context above. Cite specific documents when referencing information.""")
        ])
        
        # Generate answer
        chain = prompt | llm
        response = chain.invoke({
            "context": context[:6000],  # Increased context window
            "question": question
        })
        
        answer = response.content
        
        # Create detailed source information
        unique_sources = list(sources_dict.keys())
        source_count = len(unique_sources)
        chunk_count = len(docs)
        
        # Build source details with chunk previews (formatted for chatbot)
        sources_section = f"\n\n{'='*60}\n📚 **Sources Used** ({source_count} document(s), {chunk_count} chunk(s))\n{'='*60}\n\n"
        
        for source, source_docs in sources_dict.items():
            sources_section += f"📄 **{source}** ({len(source_docs)} chunk(s))\n"
            for idx, doc in enumerate(source_docs[:3], 1):  # Show first 3 chunks per document
                preview = doc.page_content[:150].replace("\n", " ").strip()
                if len(doc.page_content) > 150:
                    preview += "..."
                sources_section += f"   • Chunk {idx}: _{preview}_\n"
            if len(source_docs) > 3:
                sources_section += f"   • ... and {len(source_docs) - 3} more chunk(s)\n"
            sources_section += "\n"
        
        answer_with_sources = f"{answer}{sources_section}"
        
        # Update history
        history.append((question, answer_with_sources))
        
        return "", history
        
    except Exception as e:
        error_msg = f"❌ Error processing query: {str(e)}"
        history.append((question, error_msg))
        return "", history


def clear_index() -> str:
    """Clear the vectorstore index"""
    global vectorstore, retriever
    
    try:
        # Delete collection and recreate
        client = QdrantClient(location=":memory:")
        try:
            client.delete_collection(collection_name)
        except:
            pass
        
        client.create_collection(
            collection_name=collection_name,
            vectors_config=VectorParams(
                size=1536,
                distance=Distance.COSINE
            )
        )
        
        vectorstore = Qdrant(
            client=client,
            collection_name=collection_name,
            embeddings=embeddings
        )
        
        try:
            retriever = vectorstore.as_retriever(
                search_type="mmr",
                search_kwargs={"k": 10, "fetch_k": 20}
            )
        except Exception:
            retriever = vectorstore.as_retriever(
                search_kwargs={"k": 10}
            )
        
        return "✅ Index cleared successfully!"
    except Exception as e:
        return f"❌ Error clearing index: {str(e)}"


# Create Gradio interface
def create_interface():
    """Create and launch the Gradio interface"""
    
    # Initialize components
    init_status = initialize_components()
    
    # Custom CSS for better appearance
    custom_css = """
    .gradio-container {
        font-family: 'Segoe UI', Tahoma, Geneva, Verdana, sans-serif;
    }
    .main-header {
        text-align: center;
        padding: 20px;
        background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
        color: white;
        border-radius: 10px;
        margin-bottom: 20px;
    }
    .status-box {
        padding: 15px;
        background: #f0f4ff;
        border-left: 4px solid #667eea;
        border-radius: 5px;
        margin: 10px 0;
    }
    .section-header {
        background: linear-gradient(90deg, #667eea 0%, #764ba2 100%);
        color: white;
        padding: 10px 15px;
        border-radius: 5px;
        margin-bottom: 15px;
    }
    """
    
    with gr.Blocks(title="RAG Document Q&A with Qdrant", theme=gr.themes.Soft(), css=custom_css) as app:
        # Header
        gr.HTML("""
        <div class="main-header">
            <h1>📚 RAG Document Q&A System</h1>
            <p>Upload multiple documents and ask intelligent questions powered by Retrieval-Augmented Generation</p>
        </div>
        """)
        
        # Status indicator
        with gr.Row():
            status_box = gr.HTML(f"""
            <div class="status-box">
                <strong>🟢 System Status:</strong> {init_status}
            </div>
            """)
        
        with gr.Row():
            # Left column - Document Management
            with gr.Column(scale=1, min_width=400):
                gr.HTML('<div class="section-header"><h3>📤 Document Management</h3></div>')
                
                file_upload = gr.File(
                    file_count="multiple",
                    label="📎 Upload Documents",
                    file_types=[".pdf", ".doc", ".docx", ".ppt", ".pptx", ".txt"],
                    height=150
                )
                
                with gr.Row():
                    index_btn = gr.Button(
                        "📥 Index Documents", 
                        variant="primary",
                        size="lg",
                        scale=2
                    )
                    clear_btn = gr.Button(
                        "🗑️ Clear Index", 
                        variant="stop",
                        size="lg",
                        scale=1
                    )
                
                index_status = gr.Textbox(
                    label="📊 Indexing Status",
                    lines=6,
                    interactive=False,
                    show_copy_button=True
                )
                
                gr.Markdown("""
                **Supported Formats:**
                - 📄 PDF (`.pdf`)
                - 📝 Word (`.doc`, `.docx`)
                - 📊 PowerPoint (`.ppt`, `.pptx`)
                - 📋 Text (`.txt`)
                """)
            
            # Right column - Q&A Interface
            with gr.Column(scale=2, min_width=600):
                gr.HTML('<div class="section-header"><h3>💬 Ask Questions</h3></div>')
                
                chatbot = gr.Chatbot(
                    label="Conversation",
                    height=500,
                    show_copy_button=True,
                    avatar_images=(None, "🤖"),
                    bubble_full_width=False,
                    show_label=True
                )
                
                with gr.Row():
                    question_input = gr.Textbox(
                        label="",
                        placeholder="💭 Ask a question about your indexed documents... (e.g., 'What are the main topics across all documents?')",
                        lines=3,
                        scale=4
                    )
                    submit_btn = gr.Button(
                        "🚀 Submit",
                        variant="primary",
                        size="lg",
                        scale=1
                    )
                
                gr.Markdown("""
                **💡 Tips:**
                - Questions are answered using information from **all indexed documents**
                - Answers include source citations with chunk previews
                - Use MMR (Maximal Marginal Relevance) for diverse, comprehensive results
                """)
        
        # Footer with info
        gr.Markdown("""
        ---
        **🔍 RAG Features:**
        - **Multi-document search**: Queries search across all indexed documents
        - **MMR retrieval**: Ensures diverse results from different documents
        - **Source attribution**: Every answer includes detailed source citations
        - **Chunk previews**: See which document sections were used
        """)
        
        # Event handlers
        index_btn.click(
            fn=index_documents,
            inputs=[file_upload],
            outputs=[index_status]
        )
        
        clear_btn.click(
            fn=clear_index,
            outputs=[index_status]
        )
        
        submit_btn.click(
            fn=query_rag,
            inputs=[question_input, chatbot],
            outputs=[question_input, chatbot]
        )
        
        question_input.submit(
            fn=query_rag,
            inputs=[question_input, chatbot],
            outputs=[question_input, chatbot]
        )
    
    return app


if __name__ == "__main__":
    # Check for API key
    if not os.getenv("OPENAI_API_KEY"):
        print("⚠️ ERROR: OPENAI_API_KEY not found!")
        print("Please set it in your environment or create a .env file")
        exit(1)
    
    # Create and launch app
    app = create_interface()
    app.launch(
        server_name="0.0.0.0",
        server_port=7860,
        share=False
    )

